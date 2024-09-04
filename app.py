import os
from pptx import Presentation
import comtypes.client

# Define paths
pptx_template_path = 'certificate_template.pptx'
output_folder = 'certificates'
txt_file = 'names.txt'

# Create output directory if it doesn't exist
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

def generate_certificate(name):
    # Load the PowerPoint template
    prs = Presentation(pptx_template_path)
    
    # Replace placeholder with actual name while preserving style
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if '{{full name}}' in paragraph.text:
                        for run in paragraph.runs:
                            if '{{full name}}' in run.text:
                                # Replace the text while preserving the style
                                run.text = name
                                # No need to modify run.font to preserve existing styles

    # Save the modified PowerPoint to a temporary file
    temp_pptx_path = os.path.join(output_folder, f'{name}.pptx')
    prs.save(temp_pptx_path)
    
    return temp_pptx_path

def convert_to_pdf(input_file_path, output_file_path):
    # Convert file paths to Windows format
    input_file_path = os.path.abspath(input_file_path)
    output_file_path = os.path.abspath(output_file_path)
    
    # Create PowerPoint application object
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    
    # Set visibility to minimize
    powerpoint.Visible = 1
    
    # Open the PowerPoint slides
    slides = powerpoint.Presentations.Open(input_file_path)
    
    # Save as PDF (formatType = 32)
    slides.SaveAs(output_file_path, 32)
    
    # Close the slide deck
    slides.Close()
    powerpoint.Quit()

def main():
    # Read names from the text file
    with open(txt_file, 'r') as f:
        names = f.read().splitlines()
    
    for name in names:
        pptx_path = generate_certificate(name)
        pdf_path = os.path.join(output_folder, f'{name}.pdf')
        convert_to_pdf(pptx_path, pdf_path)
        
        # Remove the temporary PPTX file after saving as PDF
        os.remove(pptx_path)
        
        print(f"Certificate generated and saved as PDF for {name}")

if __name__ == '__main__':
    main()
